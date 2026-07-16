# -*- coding: utf-8 -*-
"""Builds the Tripadvisor AI Incentives deck (PPTX). On-brand, presentation-ready."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ASSET = os.path.join(HERE, "img")
os.makedirs(ASSET, exist_ok=True)

# ---------------- palette ----------------
INK   = RGBColor(0x0A, 0x15, 0x12)
INK2  = RGBColor(0x0E, 0x1F, 0x18)
INK3  = RGBColor(0x12, 0x28, 0x20)
MINT  = RGBColor(0x34, 0xE0, 0xA1)
GREEN = RGBColor(0x00, 0x97, 0x6E)
GREEND= RGBColor(0x00, 0x54, 0x3D)
PAPER = RGBColor(0xF5, 0xF8, 0xF5)
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INKT  = RGBColor(0x13, 0x22, 0x1C)   # near-black text on light
MUTED = RGBColor(0x53, 0x63, 0x5C)
MUTED_D = RGBColor(0xA7, 0xC3, 0xB8)
LINE  = RGBColor(0xDD, 0xE7, 0xE1)
POS   = RGBColor(0x14, 0xB0, 0x83)
AMBER = RGBColor(0xE0, 0x9B, 0x1E)
AMBERB= RGBColor(0xF2, 0xB5, 0x3C)
CORAL = RGBColor(0xF0, 0x5B, 0x5B)
GREY  = RGBColor(0xA7, 0xBD, 0xB2)

HEAD = "Avenir Next"
BODY = "Helvetica Neue"

SW, SH = 13.333, 7.5

# ---------------- charts (matplotlib) ----------------
def make_charts():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib import font_manager
    plt.rcParams["font.family"] = "sans-serif"
    plt.rcParams["font.sans-serif"] = ["Helvetica Neue", "Arial", "DejaVu Sans"]

    # ---- incrementality by segment ----
    labels = ["UK\nFirst","DE\nFirst","FR\nFirst","UK\nLapsed","IT\nFirst","ES\nFirst","FR\nLapsed","UK\nRepeat","DE\nRepeat"]
    vals   = [71,63,58,52,51,48,44,31,28]
    cols   = ["#14B083","#14B083","#F2B53C","#F2B53C","#F2B53C","#F2B53C","#F05B5B","#F05B5B","#F05B5B"]
    dashed = [0,0,1,0,1,0,1,0,0]
    fig, ax = plt.subplots(figsize=(8.2,3.7), dpi=200)
    bars = ax.bar(range(len(vals)), vals, color=cols, width=0.66, zorder=3)
    for i,b in enumerate(bars):
        if dashed[i]:
            b.set_alpha(0.5); b.set_hatch("////"); b.set_edgecolor(cols[i]); b.set_linewidth(1.2)
        star = "*" if dashed[i] else ""
        ax.text(b.get_x()+b.get_width()/2, vals[i]+1.4, f"{vals[i]}%{star}", ha="center", va="bottom",
                fontsize=11, fontweight="bold", color="#13221C")
    ax.set_xticks(range(len(labels))); ax.set_xticklabels(labels, fontsize=10, color="#53635C")
    ax.set_ylim(0,84); ax.set_yticks([])
    ax.axhline(50, color="#DDE7E1", lw=1, zorder=1, ls=(0,(4,4)))
    ax.text(len(vals)-0.4, 51.5, "breakeven-ish", fontsize=8.5, color="#8A9A92", ha="right")
    for s in ["top","right","left"]: ax.spines[s].set_visible(False)
    ax.spines["bottom"].set_color("#DDE7E1")
    ax.margins(x=0.01)
    plt.tight_layout(pad=0.6)
    fig.savefig(os.path.join(ASSET,"chart_incr.png"), transparent=True); plt.close(fig)

    # ---- deadweight: UK redemption vs incrementality ----
    cats = ["First","Lapsed","Repeat"]
    red  = [18,14,21]; inc = [71,52,31]
    x = range(len(cats)); w=0.36
    fig, ax = plt.subplots(figsize=(5.4,3.7), dpi=200)
    b1=ax.bar([i-w/2 for i in x], red, width=w, color="#A7BDB2", zorder=3, label="Redemption")
    b2=ax.bar([i+w/2 for i in x], inc, width=w, color="#00976E", zorder=3, label="Incrementality")
    for bs in (b1,b2):
        for b in bs:
            ax.text(b.get_x()+b.get_width()/2, b.get_height()+1.4, f"{int(b.get_height())}%",
                    ha="center", va="bottom", fontsize=10, fontweight="bold", color="#13221C")
    ax.set_xticks(list(x)); ax.set_xticklabels(cats, fontsize=11, color="#53635C")
    ax.set_ylim(0,84); ax.set_yticks([])
    for s in ["top","right","left"]: ax.spines[s].set_visible(False)
    ax.spines["bottom"].set_color("#DDE7E1")
    ax.legend(loc="upper center", ncol=2, frameon=False, fontsize=9.5, bbox_to_anchor=(0.5,1.12))
    plt.tight_layout(pad=0.6)
    fig.savefig(os.path.join(ASSET,"chart_dead.png"), transparent=True); plt.close(fig)

    # ---- cost per incremental booking = promo cost / incrementality ----
    labs=["UK\nFirst","ES\nFirst","DE\nFirst","UK\nLapsed","UK\nRepeat","DE\nRepeat"]
    cpib=[20,23,24,32,63,64]
    ccol=["#14B083","#14B083","#14B083","#F2B53C","#F05B5B","#F05B5B"]
    fig, ax = plt.subplots(figsize=(6.6,3.5), dpi=200)
    bars=ax.bar(range(len(cpib)), cpib, color=ccol, width=0.66, zorder=3)
    for i,b in enumerate(bars):
        ax.text(b.get_x()+b.get_width()/2, cpib[i]+1.2, f"{cpib[i]}", ha="center", va="bottom", fontsize=11, fontweight="bold", color="#13221C")
    ax.set_xticks(range(len(labs))); ax.set_xticklabels(labs, fontsize=10, color="#53635C")
    ax.set_ylim(0,76); ax.set_yticks([])
    for sname in ["top","right","left"]: ax.spines[sname].set_visible(False)
    ax.spines["bottom"].set_color("#DDE7E1")
    plt.tight_layout(pad=0.6)
    fig.savefig(os.path.join(ASSET,"chart_cpib.png"), transparent=True); plt.close(fig)

    # ---- value map: redemption vs incrementality, all segments ----
    pts=[("UK 1st",18,71,"#14B083"),("DE 1st",11,63,"#14B083"),("FR 1st*",12,58,"#F2B53C"),
         ("UK Laps",14,52,"#F2B53C"),("IT 1st*",10,51,"#F2B53C"),("ES 1st",9,48,"#F2B53C"),
         ("FR Laps*",9.6,44,"#F05B5B"),("UK Rpt",21,31,"#F05B5B"),("DE Rpt",14,28,"#F05B5B")]
    fig, ax = plt.subplots(figsize=(6.9,4.0), dpi=200)
    ax.axhspan(55,80,color="#14B083",alpha=0.07,zorder=0)
    ax.axhspan(20,42,color="#F05B5B",alpha=0.08,zorder=0)
    ax.axhline(50,color="#C7D3CC",lw=1,ls=(0,(4,4)),zorder=1)
    offs={"ES 1st":(2,7),"FR Laps*":(8,-12),"UK Laps":(8,2),"DE 1st":(8,2),"IT 1st*":(8,3),"UK Rpt":(9,2),"DE Rpt":(8,3)}
    for name,x,y,c in pts:
        ax.scatter(x,y,s=180,color=c,edgecolor="white",linewidth=1.6,zorder=3)
        dx,dy=offs.get(name,(8,4))
        ax.annotate(name,(x,y),xytext=(dx,dy),textcoords="offset points",fontsize=8,color="#13221C",fontweight="bold",zorder=4)
    ax.set_xlim(6.5,24.5); ax.set_ylim(20,80)
    ax.set_xlabel("Redemption  \u2014  how many take it  \u2192", fontsize=10, color="#53635C")
    ax.set_ylabel("Incrementality  \u2014  how many needed it  \u2192", fontsize=10, color="#53635C")
    ax.set_yticks([30,50,70]); ax.set_yticklabels(["30%","50%","70%"], fontsize=9, color="#8A9A92")
    ax.set_xticks([10,15,20]); ax.set_xticklabels(["10%","15%","20%"], fontsize=9, color="#8A9A92")
    ax.text(7.0,74.8,"VALUE CREATED", ha="left", fontsize=8.5, color="#14795c", fontweight="bold")
    ax.text(7.0,22.2,"VALUE DESTROYED", ha="left", fontsize=8.5, color="#c23b3b", fontweight="bold")
    for sp in ["top","right"]: ax.spines[sp].set_visible(False)
    ax.spines["left"].set_color("#DDE7E1"); ax.spines["bottom"].set_color("#DDE7E1")
    plt.tight_layout(pad=0.5)
    fig.savefig(os.path.join(ASSET,"chart_valuemap.png"), transparent=True); plt.close(fig)
    print("charts done")

# ---------------- pptx helpers ----------------
def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
def _no_line(shape):
    shape.line.fill.background()
def _line(shape, color, w=1.0):
    shape.line.color.rgb = color; shape.line.width = Pt(w)

def slide(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0, Inches(SW), Inches(SH))
    _set_fill(bg, INK if dark else PAPER); _no_line(bg)
    bg.shadow.inherit = False
    return s

def rrect(s, x,y,w,h, fill=CARD, line=None, lw=1.0, radius=0.06):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),Inches(y),Inches(w),Inches(h))
    try: shp.adjustments[0]=radius
    except Exception: pass
    if fill is None: shp.fill.background()
    else: _set_fill(shp, fill)
    if line is None: _no_line(shp)
    else: _line(shp, line, lw)
    shp.shadow.inherit=False
    return shp

def rect(s, x,y,w,h, fill, line=None, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: shp.fill.background()
    else: _set_fill(shp, fill)
    if line is None: _no_line(shp)
    else: _line(shp, line, lw)
    shp.shadow.inherit=False
    return shp

def chevron(s, x,y,w,h, fill, line=None, adj=0.18):
    shp = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x),Inches(y),Inches(w),Inches(h))
    try: shp.adjustments[0]=adj
    except Exception: pass
    if fill is None: shp.fill.background()
    else: _set_fill(shp, fill)
    if line is None: _no_line(shp)
    else: _line(shp, line, 1.0)
    shp.shadow.inherit=False
    return shp

def txt(s, x,y,w,h, runs, size=14, color=INKT, bold=False, font=BODY, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spacing=1.0, italic=False):
    """runs: string or list of (text,{overrides}) or list of strings (paragraphs)."""
    tb = s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs,str): runs=[runs]
    first=True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first=False
        p.alignment=align; p.line_spacing=spacing
        segs = para if isinstance(para, list) else [para]
        for seg in segs:
            t,ov = seg if isinstance(seg,tuple) else (seg,{})
            r=p.add_run(); r.text=t
            r.font.name=ov.get("font",font); r.font.size=Pt(ov.get("size",size))
            r.font.bold=ov.get("bold",bold); r.font.italic=ov.get("italic",italic)
            r.font.color.rgb=ov.get("color",color)
            sp=ov.get("space")
            if sp is not None:
                p.space_after=Pt(sp)
    return tb

def kicker(s, text, x, y, color=GREEN, size=12):
    txt(s, x, y, 8, 0.35, text.upper(), size=size, color=color, bold=True, font=HEAD)

def footer(s, page, dark=False):
    c = MUTED_D if dark else MUTED
    txt(s, 0.6, SH-0.5, 8, 0.3, "Adrien Enjalbert  ·  AI Systems, Personalization & Growth", size=9, color=c, font=BODY)
    txt(s, SW-1.4, SH-0.5, 0.8, 0.3, f"{page:02d}", size=9, color=c, font=HEAD, align=PP_ALIGN.RIGHT, bold=True)

def accent(s, x, y, w=0.55):
    rect(s, x, y, w, 0.055, MINT)

def note(s, text):
    s.notes_slide.notes_text_frame.text = text

def question(s, code, qtext, dark=False):
    """Prominent badge (question number) + the restated brief question, for clear mapping."""
    bw = 0.82 if len(code) <= 3 else 1.02
    rrect(s, 0.85, 0.62, bw, 0.46, MINT, radius=0.5)
    txt(s, 0.85, 0.6, bw, 0.46, code, size=14, color=INK, bold=True, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 0.85+bw+0.22, 0.6, 11.0, 0.5, [[("ANSWERING  ",{"color":(MUTED_D if dark else GREY),"size":9}),(qtext,{"color":(WHITE if dark else INKT)})]], size=12.5, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)

def build():
    prs = Presentation()
    prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    for i, fn in enumerate(ORDER, start=1):
        fn(prs, i)
    return prs

# ==== slide builders appended below ====
SLIDES = []
def register(fn):
    SLIDES.append(fn); return fn

@register
def s01(prs, pg):
    s = slide(prs, dark=True)
    # ambient bubble motif
    for (x,y,d,c) in [(11.3,0.7,1.7,INK3),(12.2,1.9,0.9,GREEND),(10.4,5.9,1.2,INK3)]:
        o=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d)); _set_fill(o,c); _no_line(o); o.shadow.inherit=False
    accent(s, 0.85, 1.55)
    kicker(s, "Interview case study  ·  International Incentives", 0.85, 1.8, MINT)
    txt(s, 0.8, 2.25, 10.6, 2.6, [
        [("Smarter incentives,", {"color":WHITE})],
        [("built to ", {"color":WHITE}), ("land", {"color":MINT}), (" in a real organisation.", {"color":WHITE})],
    ], size=44, bold=True, font=HEAD, spacing=1.02)
    txt(s, 0.85, 4.9, 9.8, 1.4, [
        "Turning a slow, meeting-heavy experiment process into a lean AI operating system that spends less, learns faster, and ships value in weeks, not quarters.",
    ], size=17, color=MUTED_D, font=BODY, spacing=1.12)
    txt(s, 0.85, 6.35, 10, 0.4, [[("AI Systems, Personalization & Growth Manager", {"color":WHITE,"bold":True}), ("   ·   prepared for the Tripadvisor interview team", {"color":MUTED_D})]], size=12.5, font=HEAD)
    footer(s, pg, dark=True)
    note(s, "One line: I read this brief less as an AI problem and more as an organisation problem. In a company this size the model is the easy part. What decides success is data you can trust, integration into the stack you already run, UX teams will actually adopt, and buy-in that survives shifting priorities. That is what this deck is built around.")

@register
def s02(prs, pg):
    s = slide(prs, dark=True)
    accent(s, 0.85, 0.75)
    kicker(s, "The reframe", 0.85, 1.0, MINT)
    txt(s, 0.8, 1.35, 11.6, 1.6, [[("In a big company, the model is the easy ", {"color":WHITE}), ("10%", {"color":MINT}), (".", {"color":WHITE})]], size=34, bold=True, font=HEAD)
    txt(s, 0.85, 2.5, 11.4, 0.7, "A recommendation engine is a weekend prototype. Getting it used, trusted, and safe across teams is the year of work. So I would spend my effort on the four things that actually block value.", size=15.5, color=MUTED_D, font=BODY, spacing=1.1)
    # 10/90 bar
    rect(s, 0.85, 3.5, 11.6, 0.62, INK3)
    rect(s, 0.85, 3.5, 11.6*0.10, 0.62, MINT)
    txt(s, 0.9, 3.56, 3, 0.5, "MODEL  ·  10%", size=11, color=INK, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 3.0, 3.56, 9.2, 0.5, "THE 90% THAT DECIDES SUCCESS", size=11, color=MUTED_D, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
    pillars = [
        ("Data & ETL", "Investigate, clean and pipe the data so results can be trusted."),
        ("Integration", "Plug into the warehouse, tracking and Braze, not a side tool."),
        ("UX & adoption", "Embed AI where teams already work, so it gets used properly."),
        ("Change management", "Small wins and coalition-building to earn resource and buy-in."),
    ]
    cw=2.78; gap=0.19; x0=0.85; y=4.48
    for i,(h,b) in enumerate(pillars):
        x=x0+i*(cw+gap)
        c=rrect(s, x, y, cw, 2.06, INK2, line=GREEND, lw=1.0)
        rect(s, x+0.28, y+0.28, 0.42, 0.06, MINT)
        txt(s, x+0.28, y+0.46, cw-0.5, 0.6, h, size=15, color=WHITE, bold=True, font=HEAD, spacing=0.98)
        txt(s, x+0.28, y+1.18, cw-0.5, 0.85, b, size=11, color=MUTED_D, font=BODY, spacing=1.05)
    footer(s, pg, dark=True)
    note(s, "This is the spine of the whole submission. I have shipped AI features before, so I know the demo is not the hard bit. In a lean org with competing priorities, the risk is a clever tool nobody adopts. These four pillars, data, integration, UX and change management, are where I would put the work, and every later slide maps back to one of them.")

@register
def s03(prs, pg):
    s = slide(prs, dark=False)
    question(s, "1.1", "Which segments create or destroy value?")
    txt(s, 0.8, 1.3, 11.8, 1.2, "One flat rule is hiding two different businesses.", size=30, bold=True, font=HEAD, color=INKT)
    txt(s, 0.85, 2.35, 6.0, 1.0, "Read by incrementality, the pattern is clear: the newer the customer, the more the discount actually causes the booking. We pay the same 15% to people who need convincing and to people who were already coming.", size=13.5, color=MUTED, font=BODY, spacing=1.1)
    s.shapes.add_picture(os.path.join(ASSET,"chart_incr.png"), Inches(0.7), Inches(3.4), width=Inches(7.8))
    # side callouts
    cx=9.35; cw=3.35
    c1=rrect(s, cx, 2.3, cw, 1.55, CARD, line=LINE); rect(s, cx, 2.3, 0.08, 1.55, POS)
    txt(s, cx+0.28, 2.5, cw-0.5, 0.4, "VALUE CREATED", size=10.5, color=POS, bold=True, font=HEAD)
    txt(s, cx+0.28, 2.86, cw-0.5, 0.9, [[("First purchase", {"bold":True,"color":INKT}), (": UK 71%, DE 63% incremental. New customers we genuinely win.", {"color":MUTED})]], size=12, font=BODY, spacing=1.05)
    c2=rrect(s, cx, 4.05, cw, 1.55, CARD, line=LINE); rect(s, cx, 4.05, 0.08, 1.55, CORAL)
    txt(s, cx+0.28, 4.25, cw-0.5, 0.4, "VALUE DESTROYED", size=10.5, color=CORAL, bold=True, font=HEAD)
    txt(s, cx+0.28, 4.61, cw-0.5, 0.9, [[("Repeat buyers", {"bold":True,"color":INKT}), (": DE 28%, UK 31%. We subsidise people who would book anyway.", {"color":MUTED})]], size=12, font=BODY, spacing=1.05)
    txt(s, cx, 5.75, cw, 0.8, [[("The metric that matters is incrementality, not redemption. ", {"bold":True,"color":INKT}), ("UK repeat redeems the most and is nearly the least incremental.", {"color":MUTED})]], size=11.5, font=BODY, spacing=1.05)
    footer(s, pg)
    note(s, "The headline insight. The flat discount averages two opposite realities. First-purchase discounts are genuinely incremental and worth protecting. Repeat discounts are largely deadweight, we are paying loyal customers to do what they would do anyway. The trap is redemption: UK repeat has the highest redemption and almost the lowest incrementality, so a busy-looking promo is quietly losing money.")

@register
def s04(prs, pg):
    s = slide(prs, dark=True)
    accent(s, 0.85, 0.75)
    kicker(s, "In practice  ·  Part 1  ·  The data foundation", 0.85, 1.0, MINT)
    txt(s, 0.8, 1.35, 11.9, 1.7, [[("\u201CIncrementality should be treated as directional.\u201D", {"color":MINT})]], size=27, bold=True, font=HEAD, spacing=1.02)
    txt(s, 0.85, 2.75, 11.4, 0.9, "That one line in the brief is the whole job. France and Italy had targeting issues in the holdout, so their numbers cannot be trusted. You cannot automate on data you cannot trust, and this is not a rare event, it is the default state of data in a large org.", size=15, color=MUTED_D, font=BODY, spacing=1.12)
    # trust pipeline chevrons
    y=4.05
    steps=[("Raw tracking","Backend events, web & app, Braze sends"),("ETL & clean-up","Dedupe, fix attribution, single source of truth"),("Validated experiment data","Holdout integrity, sample sizes, confidence"),("AI can act","Only now is automation safe")]
    n=len(steps); cw=2.9; gap=0.12; x0=0.85
    for i,(h,b) in enumerate(steps):
        x=x0+i*(cw+gap)
        col = MINT if i==n-1 else INK2
        ch=chevron(s, x, y, cw, 1.05, col, line=(None if i==n-1 else GREEND))
        tc = INK if i==n-1 else WHITE
        txt(s, x+0.36, y+0.16, cw-0.72, 0.4, h, size=12, color=tc, bold=True, font=HEAD, anchor=MSO_ANCHOR.TOP)
        txt(s, x+0.36, y+0.55, cw-0.72, 0.5, b, size=9.5, color=(INK2 if i==n-1 else MUTED_D), font=BODY, spacing=1.0)
    txt(s, 0.85, 5.55, 11.5, 0.9, [[("Where I would start: ", {"bold":True,"color":WHITE}), ("a data investigation and clean-up sprint plus a re-run of the France/Italy holdout, before a single model ships. Realistically 60\u201370% of the effort. It is unglamorous, and it is the difference between an AI that helps and one that confidently misleads.", {"color":MUTED_D})]], size=13, font=BODY, spacing=1.1)
    footer(s, pg, dark=True)
    note(s, "This is the slide I would linger on because it signals I have actually shipped in a big org. The directional flag is a symptom of shaky experiment infrastructure. If we skip the data work and point a recommendation engine at it, we get exactly the Part 3 failure. So the first deliverable is not a model, it is trustworthy data and a clean holdout. That also happens to be the cheapest way to earn credibility.")

def bullets(s, x, y, w, items, color=MUTED, size=12.5, gap=7, lead=INKT):
    runs=[]
    for it in items:
        if isinstance(it, tuple):
            head, rest = it
            runs.append([("\u25B8  ", {"color":GREEN,"bold":True}), (head, {"color":lead,"bold":True}), (rest, {"color":color})])
        else:
            runs.append([("\u25B8  ", {"color":GREEN,"bold":True}), (it, {"color":color})])
    txt(s, x, y, w, 3.5, [ [ (seg[0],seg[1]) if isinstance(seg,tuple) else seg for seg in r] for r in runs], size=size, font=BODY, spacing=1.05)
    # apply space_after
    return

@register
def s05(prs, pg):
    s = slide(prs, dark=False)
    question(s, "1.1", "Which metrics do you trust, and which would you question?")
    txt(s, 0.8, 1.3, 11.8, 1.0, "I trust the metric. I question the measurement.", size=29, bold=True, font=HEAD, color=INKT)
    txt(s, 0.85, 2.22, 11.5, 0.5, [[("The apparent contradiction resolved: ",{"bold":True,"color":INKT}),("incrementality is the right metric (it's causal). But a metric is only as good as the experiment behind it, and two of these were built on broken experiments.",{"color":MUTED})]], size=12.5, font=BODY, spacing=1.05)
    # trust
    cA=rrect(s, 0.85, 2.95, 5.75, 3.85, CARD, line=LINE); rect(s,0.85,2.95,5.75,0.09,POS)
    txt(s, 1.15, 3.25, 5.2, 0.4, "TRUST \u2014 THE METRIC TYPE", size=12, color=POS, bold=True, font=HEAD)
    txt(s, 1.15, 3.75, 5.25, 3.0, [
        [("Incrementality is causal. ",{"bold":True,"color":INKT}),("A holdout gives the counterfactual (what people do with no offer), so it isolates cause, not correlation.",{"color":MUTED})],
        [("It beats redemption & CVR lift. ",{"bold":True,"color":INKT}),("Those are observational: they count what happened, not what we caused. Incrementality nets out the free-riders.",{"color":MUTED})],
        [("Where the holdout is clean (UK, DE), ",{"bold":True,"color":INKT}),("I'd act on it, it's the one number that separates real growth from subsidised habit.",{"color":MUTED})],
    ], size=12, font=BODY, spacing=1.06)
    # question
    cB=rrect(s, 6.75, 2.95, 5.75, 3.85, CARD, line=LINE); rect(s,6.75,2.95,5.75,0.09,CORAL)
    txt(s, 7.05, 3.25, 5.2, 0.4, "QUESTION \u2014 THE MEASUREMENT", size=12, color=CORAL, bold=True, font=HEAD)
    txt(s, 7.05, 3.75, 5.25, 3.0, [
        [("FR & IT incrementality (the same metric). ",{"bold":True,"color":INKT}),("Targeting issues broke randomisation, so the estimate is biased, not just noisy. Directional only.",{"color":MUTED})],
        [("No sample size or confidence interval. ",{"bold":True,"color":INKT}),("Can't separate signal from noise; a 71% on a tiny cell isn't the same as 71% on 10k.",{"color":MUTED})],
        [("Attribution traps. ",{"bold":True,"color":INKT}),("Redemption is self-selected; last-click over-credits the offer; pooling markets risks Simpson's paradox.",{"color":MUTED})],
    ], size=12, font=BODY, spacing=1.06)
    footer(s, pg)
    note(s, "This directly answers the 'incrementality is in both columns' question. The resolution is the mark of someone who's done real experimentation: I trust incrementality as a metric TYPE because it's causal, but a causal metric is only as trustworthy as the experiment that produced it. France and Italy had broken randomisation, so their incrementality is biased, not merely uncertain. And with no sample sizes I can't attach confidence. So I trust the concept, and I interrogate the measurement. I'd also flag the classic attribution traps: redemption is self-selected, last-click over-credits, and pooling markets can flip the result via Simpson's paradox.")

@register
def s06(prs, pg):
    s = slide(prs, dark=False)
    question(s, "1.3", "Recommend the next three experiments.")
    txt(s, 0.8, 1.3, 11.9, 1.0, "Only three slots a quarter. Spend them on learning that pays.", size=27, bold=True, font=HEAD, color=INKT)
    cards=[
        ("01","Stop paying loyal customers","UK & Germany","Repeat buyers", [
            ("Offer. ","Holdout: 15% vs 0% vs a non-cash perk."),
            ("Why. ","Repeat is the clearest value-destroyer (28\u201331%). Biggest, fastest saving for Finance."),
            ("Learn. ","How much repeat revenue is truly incremental."),
        ], "Some repeat bookings may be incremental. The holdout caps exposure; watch retention, not just next booking."),
        ("02","Fixed value vs percentage","Germany","First purchase", [
            ("Offer. ","A/B \u20AC15 voucher vs 15%. Cost matched on a \u20AC102 basket."),
            ("Why. ","Test the qualitative hypothesis on the proven segment, don't assume it."),
            ("Learn. ","Whether framing lifts incrementality at equal cost."),
        ], "Framing effect may be small, but cost is matched, so the downside is a slot, not budget."),
        ("03","How shallow can we go?","Spain","First purchase", [
            ("Offer. ","Depth test: 10% vs 15%, plus a \u20AC10-off-\u20AC75 threshold."),
            ("Why. ","Marginal (48%) but not flagged, lowest basket. Cleanest place to find the floor."),
            ("Learn. ","The discount elasticity that answers Finance directly."),
        ], "Under-discounting could slow acquisition. Measure downstream CLV, not just first bookings."),
    ]
    cw=3.85; gap=0.28; x0=0.85; y=2.35
    for i,(n,h,mkt,aud,rows,risk) in enumerate(cards):
        x=x0+i*(cw+gap)
        rrect(s, x, y, cw, 3.9, CARD, line=LINE)
        rrect(s, x, y, cw, 1.02, INK, radius=0.09)
        rect(s, x, y+0.5, cw, 0.52, INK)  # square off bottom of header
        txt(s, x+0.28, y+0.16, 1.2, 0.4, n, size=13, color=MINT, bold=True, font=HEAD)
        txt(s, x+0.28, y+0.44, cw-0.5, 0.5, h, size=15, color=WHITE, bold=True, font=HEAD)
        txt(s, x+0.28, y+1.12, cw-0.5, 0.32, [[("MARKET ",{"color":GREEN,"bold":True,"size":8}),(mkt+"    ",{"color":INKT,"size":9.5}),("AUDIENCE ",{"color":GREEN,"bold":True,"size":8}),(aud,{"color":INKT,"size":9.5})]], font=HEAD)
        txt(s, x+0.28, y+1.5, cw-0.55, 1.5, [[(a,{"bold":True,"color":INKT}),(b,{"color":MUTED})] for a,b in rows], size=11, font=BODY, spacing=1.04)
        rrect(s, x+0.24, y+3.02, cw-0.48, 0.76, RGBColor(0xFC,0xF6,0xEA))
        txt(s, x+0.4, y+3.1, cw-0.72, 0.62, [[("Business risk. ",{"bold":True,"color":RGBColor(0x9A,0x6A,0x00)}),(risk,{"color":MUTED})]], size=9.5, font=BODY, spacing=1.02, anchor=MSO_ANCHOR.MIDDLE)
    # holdout note
    hn=rrect(s, 0.85, 6.34, 11.65, 0.6, RGBColor(0xFD,0xF3,0xF3), line=RGBColor(0xF3,0xC9,0xC9))
    txt(s, 1.15, 6.34, 11.1, 0.6, [[("\u270B  Deliberately not France or Italy. ",{"bold":True,"color":RGBColor(0xB0,0x34,0x34)}),("Their data is flagged. Spending a scarce slot on numbers we can't trust would be the mistake. Fix the holdout first.",{"color":MUTED})]], size=11.5, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, pg)
    note(s, "Fewer, higher-quality experiments. Each protects the base while it learns and serves Finance's goal of spending less without losing bookings. The strongest signal here is what I refuse to test: France and Italy, because acting on flagged data is how you learn the wrong thing and waste a slot.")

@register
def s07(prs, pg):
    s = slide(prs, dark=True)
    question(s, "2.1", "Design the end-to-end AI workflow.", dark=True)
    txt(s, 0.8, 1.32, 11.9, 1.0, "One operating system, built on CRISP-DM.", size=28, bold=True, font=HEAD, color=WHITE)
    txt(s, 0.85, 2.32, 11.5, 0.75, [[("The framework I used to ship a multi-agent system at Indeed Flex. ",{"color":WHITE,"bold":True}),("Business and data understanding come ",{"color":MUTED_D}),("before",{"color":MINT,"bold":True}),(" modelling, which is exactly what the ~95% of failed enterprise AI projects skip. The loop is what turns one experiment into a system.",{"color":MUTED_D})]], size=13, font=BODY, spacing=1.08)
    phases=[
        ("1","Business\nunderstanding","The \u201Cis it worth it?\u201D gate. Define the decision & metric. Sometimes the answer is a simple rule, not a model.",True),
        ("2","Data\nunderstanding","Trust the data first. The France/Italy lesson: no automation on data you can't trust.",True),
        ("3","Data\npreparation","ETL, clean, validated holdouts. The unglamorous 60\u201370%.",False),
        ("4","Modelling","Draft the experiment / offer. Simplest model that works, not the fanciest.",False),
        ("\u2713","Evaluation","Holdout, incrementality, significance + human sign-off before launch.",False),
        ("5","Deployment","Launch, monitor, write the learning brief.",False),
    ]
    n=len(phases); cw=1.92; gap=0.14; x0=0.68; y=3.35
    for i,(num,h,b,understand) in enumerate(phases):
        x=x0+i*(cw+gap)
        human = (num=="\u2713")
        col = MINT if understand else (AMBERB if human else INK2)
        ch=chevron(s, x, y, cw, 1.7, col, line=(None if (understand or human) else GREEND))
        light = understand or human
        txt(s, x+0.4, y+0.14, 0.6, 0.35, num, size=13, color=(INK if light else MINT), bold=True, font=HEAD)
        txt(s, x+0.4, y+0.44, cw-0.66, 0.55, h, size=10.5, color=(INK if light else WHITE), bold=True, font=HEAD, spacing=0.9)
        txt(s, x+0.4, y+1.02, cw-0.72, 0.62, b, size=8, color=(INK2 if light else MUTED_D), font=BODY, spacing=0.96)
    # loop + required elements
    lp=rrect(s, 0.68, 5.35, 11.97, 0.95, INK2, line=GREEND)
    txt(s, 1.0, 5.46, 11.4, 0.75, [[("\u21BB  It's a loop, not a line.  ",{"bold":True,"color":MINT}),("Deployment feeds the next Business-understanding step, so every experiment makes the system sharper. AI drafts and recalls; humans own the judgement; async gates replace the review meetings (most of the 60%).",{"color":MUTED_D})]], size=12, font=BODY, anchor=MSO_ANCHOR.MIDDLE, spacing=1.04)
    txt(s, 0.7, 6.46, 12.0, 0.35, [[("Covers all six:  ",{"bold":True,"color":WHITE}),("Inputs (hypothesis) \u00B7 Data sources (warehouse + evidence base) \u00B7 AI capabilities (recall, draft, QA, summarise) \u00B7 Human review (Evaluation gate) \u00B7 Outputs (spec, dashboard, brief) \u00B7 Feedback loop",{"color":MUTED_D})]], size=9.5, font=BODY)
    footer(s, pg, dark=True)
    note(s, "I anchored this on CRISP-DM on purpose, it's the framework I used for my Level 7 AI project at Indeed Flex. The point I'd make live: the reason most enterprise AI fails isn't the model, it's skipping business and data understanding, and never closing the feedback loop. So phases 1 and 2 are highlighted, the 'is it worth it / do we even need a model' gate sits at the front, humans own the Evaluation gate, and the whole thing is a cycle that compounds. It also cleanly covers the six required elements.")

@register
def s08(prs, pg):
    s = slide(prs, dark=False)
    accent(s, 0.85, 0.7); kicker(s, "In practice  ·  Part 2  ·  Integration", 0.85, 0.95, GREEN)
    txt(s, 0.8, 1.3, 11.9, 1.0, "Integration, not islands.", size=30, bold=True, font=HEAD, color=INKT)
    txt(s, 0.85, 2.3, 11.4, 0.7, "The value only lands if this plugs into the systems teams already run. A standalone tool gets ignored. So the AI layer sits in the middle of the existing stack, reading clean data and writing back into the tools where work happens.", size=13.5, color=MUTED, font=BODY, spacing=1.1)
    # sources
    y=3.4
    def node(x,y,w,h,title,sub,fill,line,tc,sc,ts=13):
        rrect(s,x,y,w,h,fill,line=line)
        txt(s,x+0.2,y+0.16,w-0.4,0.5,title,size=ts,color=tc,bold=True,font=HEAD,spacing=0.98)
        txt(s,x+0.2,y+h-0.62,w-0.4,0.55,sub,size=9.5,color=sc,font=BODY,spacing=1.0)
    # left: sources
    node(0.85,y,2.75,1.15,"Data warehouse","Bookings, CLV, margin",CARD,LINE,INKT,MUTED)
    node(0.85,y+1.35,2.75,1.15,"Backend tracking","Web & app events",CARD,LINE,INKT,MUTED)
    # ETL
    node(4.05,y+0.35,2.05,1.9,"ETL & clean-up","Trust layer",PAPER,GREEN,GREEND,MUTED)
    # AI center
    node(6.55,y+0.35,2.35,1.9,"AI incentive layer","Spec · recall · QA · learning",INK,None,WHITE,MUTED_D,14)
    rect(s,6.55,y+0.35+1.55,2.35,0.06,MINT)
    # right destinations
    node(9.35,y,3.15,1.15,"Braze","Lifecycle & CRM activation",CARD,LINE,INKT,MUTED)
    node(9.35,y+1.35,3.15,1.15,"Experimentation platform","Launch & holdouts",CARD,LINE,INKT,MUTED)
    # arrows
    for (ax,ay) in [(3.62,y+0.9),(6.12,y+1.2)]:
        a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(ax),Inches(ay),Inches(0.4),Inches(0.32)); _set_fill(a,GREY); _no_line(a); a.shadow.inherit=False
    a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(8.92),Inches(y+1.2),Inches(0.4),Inches(0.32)); _set_fill(a,MINT); _no_line(a); a.shadow.inherit=False
    txt(s, 0.85, 5.95, 11.6, 0.9, [[("The unsexy truth: ",{"bold":True,"color":INKT}),("most of the build is connectors, schema and permissions, not prompts. I'd scope those dependencies on day one, because in a lean org an un-owned integration is where projects quietly die.",{"color":MUTED})]], size=12.5, font=BODY, spacing=1.08)
    footer(s, pg)
    note(s, "Braze is an example of the lifecycle layer, the point is the shape: AI in the middle of the stack, not beside it. A marketer should get a recommendation inside the tool they already use to send campaigns, not in a separate dashboard they have to remember to open. I flag integration ownership early because that is the classic failure mode in big companies.")

@register
def s09(prs, pg):
    s = slide(prs, dark=True)
    accent(s, 0.85, 0.75); kicker(s, "In practice  ·  Part 2  ·  UX & adoption", 0.85, 1.0, MINT)
    txt(s, 0.8, 1.35, 11.9, 1.0, "Why this won't be \u201Ca chatbot\u201D.", size=30, bold=True, font=HEAD, color=WHITE)
    txt(s, 0.85, 2.35, 11.4, 0.7, "A generic prompt box is too vague to be used properly, so it gets abandoned. Adoption comes from AI that does one job, for one role, inside the tool they already work in.", size=14, color=MUTED_D, font=BODY, spacing=1.1)
    # left: generic chatbot (rejected)
    rrect(s, 0.85, 3.35, 3.4, 3.3, INK2, line=RGBColor(0x5A,0x2E,0x2E))
    txt(s, 1.15, 3.6, 2.9, 0.4, "\u2715  A GENERIC CHATBOT", size=12, color=CORAL, bold=True, font=HEAD)
    txt(s, 1.15, 4.1, 2.85, 2.4, [
        [("Vague, blank-box prompts",{"color":MUTED_D})],
        [("No workflow, no context",{"color":MUTED_D})],
        [("\u201CInteresting\u201D once, then ignored",{"color":MUTED_D})],
        [("Trust never builds",{"color":MUTED_D})],
    ], size=12.5, font=BODY, spacing=1.4)
    # right: embedded per role
    rrect(s, 4.55, 3.35, 7.95, 3.3, INK2, line=GREEND)
    txt(s, 4.9, 3.6, 7.4, 0.4, "\u2713  EMBEDDED, PER ROLE, PER WORKFLOW", size=12, color=MINT, bold=True, font=HEAD)
    rows=[("Analyst","gets a drafted, powered spec inside the experiment tool."),("Marketer","gets a vetted recommendation inside Braze, where they build sends."),("Leadership","gets the plain-English learning brief inside the dashboard.")]
    yy=4.15
    for r,(who,what) in enumerate(rows):
        rrect(s, 4.9, yy, 7.25, 0.72, INK, line=GREEND)
        txt(s, 5.15, yy+0.06, 2.0, 0.6, who, size=13, color=MINT, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, 7.0, yy+0.06, 5.0, 0.6, what, size=12, color=MUTED_D, font=BODY, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
        yy+=0.82
    footer(s, pg, dark=True)
    note(s, "This is the adoption argument and it is where a lot of AI initiatives fail. A chatbot puts the burden on the user to know what to ask. I'd flip it: the AI shows up already doing the job, in context, for a specific role. That is also a UX and design workload, writing the microcopy, fitting existing tools, which I'd resource deliberately rather than assume.")

@register
def s10(prs, pg):
    s = slide(prs, dark=False)
    question(s, "2.2", "Rank the AI investments. What would you build first?")
    txt(s, 0.8, 1.3, 11.9, 1.0, "Build one for quality, one for adoption.", size=28, bold=True, font=HEAD, color=INKT)
    txt(s, 0.85, 2.3, 11.4, 0.6, [[("The spec generator makes experiments better; the dashboard makes progress ",{"color":MUTED}),("visible",{"color":GREEN,"bold":True}),(". Adoption and buy-in decide ROI, not cleverness \u2014 if you don't track it and show it, it's no use. The tempting one is deliberately last.",{"color":MUTED})]], size=13, font=BODY, spacing=1.06)
    items=[
        ("1","AI experiment specification generator","Biggest bite out of the 60%. Turns a hypothesis plus history into a powered, reviewable spec.","BUILD NOW",True),
        ("2","AI dashboard summarisation","The adoption engine. Visible, plain-English results win buy-in and drive change management. If it isn't shown, it isn't used.","BUILD NOW",True),
        ("3","AI learning-brief generator","The loop-closer: turns each result into graded knowledge in the evidence base. Low risk.","SEQUENCE NEXT",False),
        ("4","AI QA assistant for launches","The safety net that would have caught the German currency error. Fast-follow.","SEQUENCE NEXT",False),
        ("5","AI promotion recommendation engine","Most powerful, highest risk. It's what failed in Part 3. Needs the data, guardrails & QA first.","EARNS IT LATER",False),
        ("6","AI customer-persona simulator","Decision-support with confirmation-bias risk. Build after the loop is trusted.","EARNS IT LATER",False),
    ]
    y=2.98; rh=0.6; gapr=0.055
    for i,(n,h,b,tag,build) in enumerate(items):
        yy=y+i*(rh+gapr)
        rrect(s, 0.85, yy, 11.65, rh, CARD if build else PAPER, line=(GREEN if build else LINE))
        # rank num
        txt(s, 1.05, yy, 0.7, rh, n, size=22, color=(GREEN if build else GREY), bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        txt(s, 1.9, yy+0.05, 4.7, rh-0.1, h, size=13, color=INKT, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE, spacing=0.98)
        txt(s, 6.75, yy+0.05, 4.2, rh-0.1, b, size=10.5, color=MUTED, font=BODY, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
        pill_c = GREEN if build else (AMBER if "NEXT" in tag else GREY)
        txt(s, 10.95, yy, 1.5, rh, tag, size=9, color=pill_c, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    footer(s, pg)
    note(s, "My re-ranked call: build one thing for quality (the spec generator) and one for adoption (the dashboard). I moved the dashboard up deliberately, because in a big org a system nobody can see the value of doesn't get resourced, visibility is how you earn change management and buy-in. The learning brief and QA follow fast. The recommendation engine is last on purpose: it's the flashiest and it's exactly what blew up in Part 3, so it only earns autonomy once the data, guardrails and QA exist.")

@register
def s11(prs, pg):
    s = slide(prs, dark=True)
    question(s, "P3", "Give 5+ reasons the AI recommendation failed.", dark=True)
    txt(s, 0.8, 1.35, 12.0, 1.1, [[("The engine swapped 15% for a \u20AC10 offer in Germany. It flopped. Why?",{"color":WHITE})]], size=24, bold=True, font=HEAD, spacing=1.02)
    txt(s, 0.85, 2.55, 11.6, 0.6, [[("First suspect: a flat \u20AC10 is ",{"color":MUTED_D}),("less",{"color":MINT,"bold":True}),(" than the \u20AC15.30 that 15% gave on a \u20AC102 basket, and it was even quoted in pounds. It under-priced the offer.",{"color":MUTED_D})]], size=13.5, font=BODY, spacing=1.08)
    reasons=[
        ("Currency & basket mismatch","Recompute effective discount vs the 15% baseline; check the currency field (\u00A3 vs \u20AC)."),
        ("Applied to the wrong segment","Break the live result by segment; confirm who received it vs where the signal came from."),
        ("Soft signal as a hard rule","Pull the qual study's sample size, date and segment; test where the preference holds."),
        ("Correlation, not incrementality","Audit the objective & features; re-rank the pick on holdout incrementality only."),
        ("No clean baseline","Confirm a holdout existed; check seasonality, competitor promos and tracking breaks."),
        ("Confident on thin data","Check training-data coverage for fixed-value offers and whether a confidence score was ignored."),
    ]
    cw=3.85; gap=0.28; x0=0.85; y=3.28; ch=1.4
    for i,(h,b) in enumerate(reasons):
        col=i%3; row=i//3
        x=x0+col*(cw+gap); yy=y+row*(ch+0.18)
        rrect(s, x, yy, cw, ch, INK2, line=GREEND)
        txt(s, x+0.28, yy+0.18, 0.6, 0.4, str(i+1), size=16, color=MINT, bold=True, font=HEAD)
        txt(s, x+0.82, yy+0.16, cw-1.05, 0.5, h, size=12.5, color=WHITE, bold=True, font=HEAD, spacing=0.96)
        txt(s, x+0.82, yy+0.62, cw-1.05, 0.72, [[("Investigate \u2192 ",{"bold":True,"color":MINT}),(b,{"color":MUTED_D})]], size=9.7, font=BODY, spacing=1.02)
    txt(s, 0.85, 6.5, 11.6, 0.6, [[("The lesson: ",{"bold":True,"color":MINT}),("guardrails before autonomy. The QA assistant and an incrementality-first knowledge base would have caught this. That's why they rank above the engine.",{"color":MUTED_D})]], size=11.5, font=BODY, spacing=1.04)
    footer(s, pg, dark=True)
    note(s, "I'd investigate each with a specific check, but the pattern matters more than the list: the engine failed quietly because the offer was mis-priced, aimed at the wrong segment, and optimised on correlation rather than incrementality. This is the evidence for the build order on the previous slide.")

@register
def s12(prs, pg):
    s = slide(prs, dark=False)
    accent(s, 0.85, 0.7); kicker(s, "In practice  ·  Part 4  ·  Change management", 0.85, 0.95, GREEN)
    txt(s, 0.8, 1.3, 11.9, 1.0, "Land it with small wins, not a big-bang rollout.", size=28, bold=True, font=HEAD, color=INKT)
    txt(s, 0.85, 2.3, 11.4, 0.65, "A large, resource-constrained org with shifting priorities rewards proof over promises. So I'd earn trust and budget in stages, keeping dependencies small at each step.", size=13.5, color=MUTED, font=BODY, spacing=1.08)
    phases=[
        ("Weeks 0\u20134","Prove the data","Data clean-up + re-run the France/Italy holdout. First visible win: numbers people trust.","No engineering ask yet"),
        ("Weeks 4\u201310","Prove the loop","Ship the spec + learning-brief generator with one friendly team. One experiment, start to finish.","1 engineer, 1 team"),
        ("Weeks 10\u201320","Prove the value","Show time-to-launch down and a saving from the repeat-customer holdout. Take it to Finance.","Exec sponsor, scorecard"),
        ("Quarter 2+","Expand","Add QA, then earn the recommendation engine. Widen to more markets and teams.","Wider resource, on evidence"),
    ]
    cw=2.9; gap=0.18; x0=0.85; y=3.2
    for i,(wk,h,b,ask) in enumerate(phases):
        x=x0+i*(cw+gap)
        rrect(s, x, y, cw, 3.15, CARD, line=LINE)
        rect(s, x, y, cw, 0.09, MINT)
        txt(s, x+0.24, y+0.28, cw-0.4, 0.35, wk.upper(), size=10.5, color=GREEN, bold=True, font=HEAD)
        txt(s, x+0.24, y+0.66, cw-0.45, 0.5, h, size=15, color=INKT, bold=True, font=HEAD, spacing=0.98)
        txt(s, x+0.24, y+1.25, cw-0.45, 1.4, b, size=11, color=MUTED, font=BODY, spacing=1.06)
        rrect(s, x+0.24, y+2.55, cw-0.48, 0.45, PAPER)
        txt(s, x+0.24, y+2.55, cw-0.48, 0.45, ask, size=9.5, color=GREEND, bold=True, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, pg)
    note(s, "This is the part that most AI plans skip. In a lean org you do not get a blank cheque, you get a chance to prove something small. Each phase produces a win that unlocks the next and keeps dependencies minimal. The data clean-up is first precisely because it needs no engineering and instantly builds credibility with Finance and analysts.")

@register
def s13(prs, pg):
    s = slide(prs, dark=True)
    question(s, "P4", "Align Finance and the Incentives Lead.", dark=True)
    txt(s, 0.8, 1.35, 12.0, 1.0, "Finance wants proof. The Lead wants the waste gone. Both are right.", size=24, bold=True, font=HEAD, color=WHITE)
    # two quotes
    q1=rrect(s, 0.85, 2.55, 5.75, 1.5, INK2, line=RGBColor(0x5A,0x2E,0x2E))
    txt(s, 1.15, 2.72, 5.2, 0.35, "FINANCE", size=10.5, color=CORAL, bold=True, font=HEAD)
    txt(s, 1.15, 3.08, 5.25, 0.95, [[("\u201CIt's delivering bookings. I'm uncomfortable reducing discounts before we have stronger evidence.\u201D",{"italic":True,"color":RGBColor(0xE6,0xEF,0xEA)})]], size=12.5, font=BODY, spacing=1.05)
    q2=rrect(s, 6.75, 2.55, 5.75, 1.5, INK2, line=GREEND)
    txt(s, 7.05, 2.72, 5.2, 0.35, "INCENTIVES LEAD  ·  PRIVATELY", size=10.5, color=MINT, bold=True, font=HEAD)
    txt(s, 7.05, 3.08, 5.25, 0.95, [[("\u201CWe're wasting promotional budget on customers who would have booked anyway.\u201D",{"italic":True,"color":RGBColor(0xE6,0xEF,0xEA)})]], size=12.5, font=BODY, spacing=1.05)
    # resolution: data-driven, sampled, both in the loop
    rrect(s, 0.85, 4.3, 11.65, 2.1, INK2, line=GREEND)
    txt(s, 1.15, 4.46, 11.1, 0.4, "MY MOVE  ·  TEST ON A SAMPLE  \u00B7  DECIDE ON THE NUMBERS  \u00B7  REPORT TO BOTH", size=11, color=MINT, bold=True, font=HEAD)
    steps=[
        ("Sample, don't switch.","A holdout on repeat buyers in 1\u20132 markets caps the downside. The promo runs everywhere else, so Finance's bookings stay protected."),
        ("Pre-agree the decision rule.","We write down now what result cuts, keeps or scales the discount, so evidence decides, not opinion or seniority."),
        ("One live scorecard, both in the loop.","Finance and the Lead see the same weekly read of spend, bookings and incrementality. Visibility replaces the debate."),
    ]
    yy=4.86
    for h,b in steps:
        txt(s, 1.15, yy, 11.2, 0.5, [[("\u25B8  ",{"color":MINT,"bold":True}),(h+" ",{"bold":True,"color":WHITE}),(b,{"color":MUTED_D})]], size=11.5, font=BODY, spacing=1.02)
        yy+=0.5
    footer(s, pg, dark=True)
    note(s, "This is deliberately experiment-led and low-drama. I refuse the false cut-vs-keep choice. Instead: sample the risk with a holdout so we never bet the whole programme; pre-agree the decision rule up front so the result decides rather than whoever's most senior; and put both stakeholders on one live scorecard so we win by numbers, not opinion. It hits everything they're looking for, recognising uncertainty, reducing risk while learning, and communicating with both a finance and a commercial audience.")

@register
def s14(prs, pg):
    s = slide(prs, dark=False)
    question(s, "2.5", "How would you measure success at six months?")
    txt(s, 0.8, 1.3, 11.9, 1.0, "Different metrics matter at different stages.", size=28, bold=True, font=HEAD, color=INKT)
    txt(s, 0.85, 2.3, 11.5, 0.6, [[("Time saved",{"bold":True,"color":INKT}),(" is the hypothesis that wins the go-ahead; ",{"color":MUTED}),("quality",{"bold":True,"color":INKT}),(" proves it works; ",{"color":MUTED}),("adoption",{"bold":True,"color":INKT}),(" proves it landed; ",{"color":MUTED}),("ROI",{"bold":True,"color":INKT}),(" proves it was worth running.",{"color":MUTED})]], size=13, font=BODY, spacing=1.06)
    stages=[
        ("1","BEFORE","the hypothesis","Time saved",[ "Analyst hours per experiment reclaimed","Idea-to-launch time (\u221260% target)","The business case that wins the go-ahead"]),
        ("2","DURING BUILD","does it work?","System quality",[ "Spec accuracy & acceptance rate","Experiment quality (reach significance)","QA catch rate; recommendation hit-rate"]),
        ("3","AT LAUNCH","did it land?","Adoption",[ "Active users & experiments run through it","Rests on UX + enablement + change mgmt","Then: stakeholder satisfaction"]),
        ("4","6 MONTHS","was it worth it?","Revenue & ROI",[ "Net incremental revenue per \u00A3 of promo","Spend cut while bookings hold","ROI net of run-cost (AI isn't free to run)"]),
    ]
    cw=2.75; gap=0.2; x0=0.85; y=3.15; hh=0.98
    rect(s, 1.0, y+hh/2, 11.3, 0.03, LINE)  # timeline rail
    for i,(num,ph,q,fam,rows) in enumerate(stages):
        x=x0+i*(cw+gap)
        rrect(s, x, y, cw, 3.4, CARD, line=LINE)
        rrect(s, x, y, cw, hh, INK, radius=0.09); rect(s, x, y+hh/2, cw, hh/2, INK)
        txt(s, x+0.24, y+0.14, cw-0.4, 0.3, [[(num+"  ",{"color":MINT,"bold":True}),(ph,{"color":WHITE,"bold":True})]], size=11, font=HEAD)
        txt(s, x+0.24, y+0.42, cw-0.4, 0.3, q, size=10, color=MUTED_D, font=BODY, italic=True)
        txt(s, x+0.24, y+0.66, cw-0.4, 0.3, fam, size=13.5, color=MINT, bold=True, font=HEAD)
        txt(s, x+0.24, y+1.16, cw-0.44, 2.1, [[("\u25B8  ",{"color":GREEN,"bold":True}),(r,{"color":MUTED})] for r in rows], size=10.5, font=BODY, spacing=1.28)
    txt(s, 0.85, 6.72, 11.6, 0.4, [[("Adoption is the real leading indicator: ",{"bold":True,"color":INKT}),("a technically great system nobody uses returns nothing, which is why UX and enablement sit inside the metric, not beside it.",{"color":MUTED})]], size=11, font=BODY)
    footer(s, pg)
    note(s, "I'd frame success as a maturity curve, not one number. Before we build, time saved is the hypothesis that justifies the investment, that's the business case. During the build I watch system quality: accuracy, experiment quality, QA catch rate. At launch the metric that matters is adoption, and adoption depends as much on UX, enablement and how we roll it out as on the model itself, which then shows up as stakeholder satisfaction. Only at six months do I judge revenue and ROI, and I'd net ROI against run-cost because these systems are genuinely expensive to operate. This mirrors how I measured my Indeed Flex project: adoption (queries/week) and ROI, not just accuracy.")

@register
def s15(prs, pg):
    s = slide(prs, dark=True)
    for (x,y,d,c) in [(11.2,0.6,1.8,INK3),(12.4,4.9,1.1,GREEND),(10.2,6.1,0.9,INK3)]:
        o=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d)); _set_fill(o,c); _no_line(o); o.shadow.inherit=False
    accent(s, 0.85, 1.35); kicker(s, "The one idea", 0.85, 1.6, MINT)
    txt(s, 0.8, 2.0, 11.2, 2.3, [
        [("The flat 15% isn't a discount strategy.",{"color":WHITE})],
        [("It's an ",{"color":WHITE}),("un-tested assumption",{"color":MINT}),(" running across nine segments.",{"color":WHITE})],
    ], size=30, bold=True, font=HEAD, spacing=1.05)
    txt(s, 0.85, 4.35, 10.8, 0.9, "My whole submission is one move: make the assumption testable, build the loop that learns faster than we can argue, and land it with the data, integration, UX and buy-in that make it stick.", size=15, color=MUTED_D, font=BODY, spacing=1.12)
    # what I'd need
    rrect(s, 0.85, 5.5, 11.65, 1.05, INK2, line=GREEND)
    txt(s, 1.15, 5.62, 3.0, 0.4, "TO START, I'D NEED", size=11, color=MINT, bold=True, font=HEAD)
    needs=["Access to warehouse & experiment data","One engineer for the loop","An exec sponsor + a friendly first team"]
    xx=1.15
    txt(s, 1.15, 5.98, 11.2, 0.5, [[ ("\u2713 "+needs[0]+"      ",{"color":WHITE}),("\u2713 "+needs[1]+"      ",{"color":WHITE}),("\u2713 "+needs[2],{"color":WHITE}) ]], size=12, font=BODY)
    footer(s, pg, dark=True)
    note(s, "Close on the single idea and make the ask concrete. I'm not pitching a moonshot, I'm pitching a testable assumption plus a learning loop, delivered in stages. The three things I'd need are small and specific, which is deliberate: it signals I know how to start in a constrained org and build from a first win.")

@register
def s_info(prs, pg):
    s = slide(prs, dark=False)
    question(s, "1.2", "What additional information would you request?")
    txt(s, 0.8, 1.3, 11.9, 1.0, "Six things I'd ask for, and how each changes the call.", size=28, bold=True, font=HEAD, color=INKT)
    txt(s, 0.85, 2.3, 11.4, 0.5, "I wouldn't finalise on this table alone. Each of these moves a specific decision, not just adds detail. If a data point can't change what I do, I don't wait for it.", size=13, color=MUTED, font=BODY, spacing=1.05)
    cards=[
        ("Customer lifetime value","First-booking value is a snapshot; LTV is the prize.","A high-LTV first customer justifies a deeper discount, it's LTV:CAC, not first-order profit. Could flip 'marginal' ES/IT/FR into worth winning."),
        ("Contribution margin","Promo cost is clear; the profit it eats isn't.","Sets the real breakeven discount depth, and whether a \u201Cnegative\u201D verdict is mild or severe."),
        ("Sample size & confidence","No n or intervals are shown.","Decides whether I roll out, retest, or ignore a row. Urgent given the France/Italy flags."),
        ("Traffic source","A coupon-site redemption is low-incrementality by design.","Changes targeting: I may exclude deal-seeker channels rather than re-discount the same people."),
        ("Device & app usage","App / logged-in users skew loyal and repeat.","Changes eligibility: suppress the offer for high-frequency app users who convert without it."),
        ("Holdout integrity","The France/Italy targeting issue is the biggest unknown.","Decides whether those markets earn an experiment slot, or a measurement fix first."),
    ]
    cw=3.75; gap=0.2; x0=0.85; ys=[2.95,5.0]; rh=1.9
    for i,(h,sub,why) in enumerate(cards):
        col=i%3; row=i//3
        x=x0+col*(cw+gap); y=ys[row]
        rrect(s, x, y, cw, rh, CARD, line=LINE)
        rect(s, x, y, 0.08, rh, MINT)
        txt(s, x+0.28, y+0.2, cw-0.5, 0.4, h, size=14.5, color=INKT, bold=True, font=HEAD)
        txt(s, x+0.28, y+0.6, cw-0.5, 0.5, sub, size=11, color=MUTED, font=BODY, italic=True, spacing=1.02)
        txt(s, x+0.28, y+1.12, cw-0.5, 0.75, [[("Changes: ",{"bold":True,"color":GREEND}),(why,{"color":MUTED})]], size=11, font=BODY, spacing=1.04)
    footer(s, pg)
    note(s, "This is the 'ask good questions before recommending' signal the brief rewards. The discipline is that every request is tied to a decision it would flip, not data for its own sake. CLV and margin decide whether marginal acquisition segments are worth it; sample size and holdout integrity decide what I'm even allowed to act on; traffic and device change targeting and eligibility.")

@register
def s_know(prs, pg):
    s = slide(prs, dark=True)
    question(s, "2.3", "Design the knowledge-management system.", dark=True)
    txt(s, 0.8, 1.32, 11.9, 0.9, "Not a document store. A learning loop that grades its own evidence.", size=23, bold=True, font=HEAD, color=WHITE)
    txt(s, 0.85, 2.26, 11.6, 0.55, [[("Simplest version that works: ",{"bold":True,"color":MINT}),("store every experiment as a graded record, keep one living belief per segment, and let two agents read and write it. The evidence base in the middle is the only thing that compounds.",{"color":MUTED_D})]], size=12.5, font=BODY, spacing=1.04)
    # three layers
    bw=3.42; ay=2.92; ah=2.14
    xs=[0.85, 4.82, 8.79]
    layers=[
        ("raw/","the ledger", GREEND,
         "Every experiment as a structured record: holdout, sample size, incrementality + confidence, cost, decision, caveats. Append-only; each record carries an evidence grade.", None),
        ("context/","the evidence base  \u00B7  the asset", MINT,
         "One living belief per market \u00D7 segment \u00D7 offer: current effect, uncertainty, grade and last-verified date. This is what compounds.",
         "A holdout \u203A B flawed \u203A C observed \u203A D hunch"),
        ("views/","auto-generated, always in sync", GREEND,
         "Dashboard, learning note and draft spec, rendered from the evidence base, never hand-written. Cheap, disposable, always current.", None),
    ]
    for i,(nm,tag,ln,body,legend) in enumerate(layers):
        x=xs[i]
        rrect(s, x, ay, bw, ah, INK2, line=ln, lw=(1.7 if ln==MINT else 1.0))
        txt(s, x+0.24, ay+0.16, bw-0.4, 0.4, nm, size=17, color=(MINT if ln==MINT else WHITE), bold=True, font="IBM Plex Mono")
        txt(s, x+0.24, ay+0.56, bw-0.4, 0.3, tag.upper(), size=9, color=(MINT if ln==MINT else MUTED_D), bold=True, font=HEAD)
        txt(s, x+0.24, ay+0.92, bw-0.45, 1.05, body, size=10.3, color=MUTED_D, font=BODY, spacing=1.03)
        if legend:
            rrect(s, x+0.22, ay+ah-0.44, bw-0.44, 0.32, INK, radius=0.4)
            txt(s, x+0.22, ay+ah-0.46, bw-0.44, 0.32, legend, size=8.6, color=MINT, bold=True, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # arrows + agents between layers
    for (ax,lab,sub) in [(4.28,"Analytics agent","grade \u00B7 update \u00B7 flag"),(8.25,"Spec agent","retrieve \u00B7 draft")]:
        a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(ax),Inches(ay+0.72),Inches(0.5),Inches(0.4)); _set_fill(a,MINT); _no_line(a); a.shadow.inherit=False
        txt(s, ax-0.42, ay+1.18, 1.35, 0.5, [[(lab,{"bold":True,"color":MINT,"size":8.5})],[(sub,{"color":MUTED_D,"size":7.3})]], font=HEAD, align=PP_ALIGN.CENTER, spacing=0.9)
    # THE RULE (flagship, ties to Part 3)
    rrect(s, 0.85, 5.3, 11.65, 0.58, INK2, line=MINT)
    txt(s, 1.05, 5.28, 11.3, 0.58, [[("\u26A1  The rule that would have caught Part 3:  ",{"bold":True,"color":MINT}),("a grade-D hunch can never overrule a grade-A holdout. Evidence quality decides, not recency or seniority.",{"color":WHITE})]], size=12, font=BODY, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
    # human validation + measurement
    txt(s, 0.85, 6.02, 11.7, 0.45, [[("Humans validate:  ",{"bold":True,"color":WHITE}),("a new belief lands as ",{"color":MUTED_D}),("candidate",{"color":MINT,"bold":True}),(" \u2192 an analyst signs off and grades it \u2192 ",{"color":MUTED_D}),("confirmed",{"color":MINT,"bold":True}),(". Contradictions are flagged for review, never silently overwritten.",{"color":MUTED_D})]], size=10.3, font=BODY, spacing=1.02)
    txt(s, 0.85, 6.5, 11.7, 0.4, [[("Measured by decision quality, not pages:  ",{"bold":True,"color":WHITE}),("did retrieved evidence change a call, was it right, is it still fresh. Plain files you own, not locked in a vendor.",{"color":MUTED_D})]], size=9.7, font=BODY, spacing=1.02)
    footer(s, pg, dark=True)
    note(s, "The upgrade from a plain 'second brain' is that this system grades its own evidence. Every experiment is stored as a structured record with an evidence grade: A for a clean holdout down to D for a qualitative hunch. The middle layer, the evidence base, is the only thing that compounds, one living belief per market-segment-offer with its effect size, uncertainty and how fresh it is. Two agents work it: an Analytics agent closes the loop (grade a finished experiment, update the belief, flag contradictions, regenerate the dashboard and note), and a Spec agent opens the next one (pull the graded priors, draft a reviewable spec). The one rule that matters, and the thing that would have prevented the Part 3 German failure, is that a grade-D hunch can never overrule a grade-A holdout. Humans validate through a candidate-to-confirmed sign-off, and success is whether retrieved evidence actually changed a decision, not how many pages it holds.")

@register
def s_persona(prs, pg):
    s = slide(prs, dark=False)
    question(s, "2.4", "How would you build AI customer personas?")
    txt(s, 0.8, 1.3, 11.9, 0.95, "Five personas, built from behaviour, not stereotypes.", size=27, bold=True, font=HEAD, color=INKT)
    txt(s, 0.85, 2.24, 11.6, 0.55, [[("What ",{"bold":True,"color":GREEN}),("a data-built profile that predicts how a market reacts.  ",{"color":MUTED}),("How ",{"bold":True,"color":GREEN}),("from booking + holdout behaviour, not clich\u00E9s.  ",{"color":MUTED}),("Why ",{"bold":True,"color":GREEN}),("to pressure-test an offer before spending a scarce slot, always an input to a real test, never the verdict.",{"color":MUTED})]], size=12, font=BODY, spacing=1.04)
    ppl=[
        ("UK",POS,"Deal-aware, high basket","Very incremental first (71%), wasteful repeat (31%).","Where to stop discounting loyal buyers first."),
        ("FR",AMBER,"Read with caution","Targeting issue makes 58% directional, not fact.","What to fix before we test here."),
        ("DE",POS,"Framing-sensitive","Strong first-purchase; prefers fixed value (qual).","Pre-test fixed vs % at equal cost."),
        ("IT",AMBER,"Read with caution","Flagged; low basket (\u20AC79) shifts fixed-offer maths.","Sequencing: measure first, spend later."),
        ("ES",AMBER,"Clean but marginal","48% on the lowest basket (\u20AC75).","How shallow a discount still acquires."),
    ]
    cw=2.28; gap=0.15; x0=0.85; y=2.95; rh=1.95
    for i,(code,cc,trait,data,use) in enumerate(ppl):
        x=x0+i*(cw+gap)
        rrect(s, x, y, cw, rh, CARD, line=LINE)
        cir=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.24),Inches(y+0.22),Inches(0.5),Inches(0.5)); _set_fill(cir, INK); _no_line(cir); cir.shadow.inherit=False
        txt(s, x+0.24, y+0.22, 0.5, 0.5, code, size=12, color=MINT, bold=True, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x+0.85, y+0.26, cw-0.95, 0.45, trait, size=11.5, color=INKT, bold=True, font=HEAD, spacing=0.95)
        txt(s, x+0.24, y+0.85, cw-0.42, 0.6, data, size=10, color=MUTED, font=BODY, spacing=1.02)
        txt(s, x+0.24, y+1.45, cw-0.42, 0.45, [[("Use for: ",{"bold":True,"color":GREEND}),(use,{"color":MUTED})]], size=9.5, font=BODY, spacing=1.0)
    # two guardrail panels
    py=5.1; pw=5.75; ph=1.55
    rrect(s, 0.85, py, pw, ph, CARD, line=LINE); rect(s,0.85,py,pw,0.08,POS)
    txt(s, 1.15, py+0.16, pw-0.5, 0.35, "HOW I KEEP THEM HONEST", size=11, color=POS, bold=True, font=HEAD)
    txt(s, 1.15, py+0.55, pw-0.55, 0.95, [
        [("\u25B8  ",{"color":GREEN,"bold":True}),("Every trait traces to a booking or holdout number, not a stereotype.",{"color":MUTED})],
        [("\u25B8  ",{"color":GREEN,"bold":True}),("Log persona predictions vs live results and recalibrate. A persona that always agrees is a red flag.",{"color":MUTED})],
    ], size=10.5, font=BODY, spacing=1.06)
    rrect(s, 6.75, py, pw, ph, CARD, line=LINE); rect(s,6.75,py,pw,0.08,CORAL)
    txt(s, 7.05, py+0.16, pw-0.5, 0.35, "WHEN I IGNORE IT", size=11, color=CORAL, bold=True, font=HEAD)
    txt(s, 7.05, py+0.55, pw-0.55, 0.95, [
        [("\u25B8  ",{"color":CORAL,"bold":True}),("When a live holdout contradicts it. Real data always wins.",{"color":MUTED})],
        [("\u25B8  ",{"color":CORAL,"bold":True}),("On flagged markets, novel offers, or high-stakes one-way doors. It prompts a test, it doesn't replace one.",{"color":MUTED})],
    ], size=10.5, font=BODY, spacing=1.06)
    footer(s, pg)
    note(s, "Personas are the part most likely to be misused, so I'm deliberate: they generate hypotheses, they don't make decisions. Each is grounded in the segment data, France and Italy are explicitly low-trust because of the flag, and I bake in bias controls, tracking prediction vs reality, and clear rules for when to ignore it. That balance of automation and oversight is exactly what the brief asks for.")

@register
def s_map(prs, pg):
    s = slide(prs, dark=True)
    accent(s, 0.85, 0.75); kicker(s, "How to read this", 0.85, 1.0, MINT)
    txt(s, 0.8, 1.35, 11.9, 1.0, "Every question in the brief, answered and labelled.", size=28, bold=True, font=HEAD, color=WHITE)
    txt(s, 0.85, 2.4, 11.4, 0.5, "Each slide carries a badge with the question it answers. Four green \u201Cin practice\u201D slides sit next to the part they support, showing how it actually lands in a real org.", size=13, color=MUTED_D, font=BODY, spacing=1.05)
    cols=[
        ("PART 1","Data reasoning",[("1.1","Segments + deeper analysis"),("1.1","Metrics I trust vs question"),("1.2","Extra data I'd request"),("1.3","The next three experiments")]),
        ("PART 2","AI systems design",[("2.1","End-to-end workflow"),("2.2","Which two to build first"),("2.3","Knowledge management"),("2.4","AI personas"),("2.5","Measuring success")]),
        ("PART 3","AI debugging",[("P3","5+ reasons the rec failed"),("","+ how I'd investigate each")]),
        ("PART 4","Stakeholders",[("P4","Aligning Finance & the Lead"),("","~300-word response"),("+","In practice: change mgmt")]),
    ]
    cw=2.9; gap=0.18; x0=0.85; y=3.3; rh=3.2
    for i,(pt,ptt,rows) in enumerate(cols):
        x=x0+i*(cw+gap)
        rrect(s, x, y, cw, rh, INK2, line=GREEND)
        txt(s, x+0.26, y+0.24, cw-0.4, 0.35, pt, size=12, color=MINT, bold=True, font=HEAD)
        txt(s, x+0.26, y+0.58, cw-0.4, 0.4, ptt, size=15, color=WHITE, bold=True, font=HEAD)
        yy=y+1.15
        for code,label in rows:
            if code:
                rrect(s, x+0.26, yy, 0.5, 0.32, (MINT if code not in ("","+") else INK2), radius=0.4)
                if code not in ("","+"):
                    txt(s, x+0.26, yy-0.02, 0.5, 0.32, code, size=9.5, color=INK, bold=True, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                elif code=="+":
                    txt(s, x+0.26, yy-0.02, 0.5, 0.32, "\u2192", size=11, color=MINT, bold=True, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            txt(s, x+0.86, yy-0.04, cw-1.05, 0.42, label, size=10.3, color=MUTED_D, font=BODY, anchor=MSO_ANCHOR.MIDDLE, spacing=0.98)
            yy+=0.46
    footer(s, pg, dark=True)
    note(s, "Quick orientation so you can follow the mapping: the deck answers the brief in its own order, each slide is badged with its question number, and I close with the four things that make this work in a large organisation. Feel free to jump to any question.")

@register
def s_analytics(prs, pg):
    s = slide(prs, dark=False)
    question(s, "1.1", "Going deeper: what advanced analysis reveals.")
    txt(s, 0.8, 1.3, 11.9, 1.0, "Redemption is a vanity metric. Map it against value.", size=26, bold=True, font=HEAD, color=INKT)
    txt(s, 0.85, 2.28, 11.6, 0.55, [[("Plot every segment two ways: ",{"bold":True,"color":INKT}),("how many take the discount (redemption) vs how many actually needed it (incrementality). Winners sit high; the money leaks bottom-right.",{"color":MUTED})]], size=12.5, font=BODY, spacing=1.05)
    s.shapes.add_picture(os.path.join(ASSET,"chart_valuemap.png"), Inches(0.55), Inches(2.95), width=Inches(7.05))
    # insight + maths panel
    fx=7.9; fw=4.6
    rrect(s, fx, 2.95, fw, 1.75, CARD, line=LINE); rect(s, fx, 2.95, 0.08, 1.75, CORAL)
    txt(s, fx+0.28, 3.14, fw-0.5, 0.4, "THE ANGLE", size=10.5, color=CORAL, bold=True, font=HEAD)
    txt(s, fx+0.28, 3.54, fw-0.55, 1.1, [[("High redemption looks like success. ",{"bold":True,"color":INKT}),("But UK Repeat redeems the MOST (21%) and is nearly the least incremental (31%). That's spend, not growth, hiding as a busy promo.",{"color":MUTED})]], size=11.5, font=BODY, spacing=1.08)
    rrect(s, fx, 4.9, fw, 1.5, PAPER, line=LINE)
    txt(s, fx+0.28, 5.08, fw-0.5, 0.4, "THE MATHS BEHIND IT", size=10.5, color=GREEN, bold=True, font=HEAD)
    txt(s, fx+0.28, 5.48, fw-0.55, 0.95, [
        [("Cost / incremental booking",{"bold":True,"color":INKT}),(" = promo \u00F7 incrementality \u2192 repeat costs ~3\u00D7 more.",{"color":MUTED})],
        [("Deadweight",{"bold":True,"color":INKT}),(" = promo \u00D7 (1\u2212incr) \u2192 ~70% of repeat spend wasted.",{"color":MUTED})],
    ], size=10.5, font=BODY, spacing=1.12)
    txt(s, 7.9, 6.5, 4.6, 0.4, [[("Next, with more data: ",{"bold":True,"color":GREEND}),("iROAS & LTV:CAC per segment, CUPED, uplift modelling.",{"color":MUTED})]], size=9.5, font=BODY, spacing=1.0)
    footer(s, pg)
    note(s, "The angle here: redemption is a vanity metric, so I map every segment on two axes, redemption (how many take it) against incrementality (how many actually needed it). The winners sit high regardless of redemption; the value leaks in the bottom-right, where UK Repeat redeems the most and needs it least. Then I back it with the maths: cost per incremental booking (promo divided by incrementality) shows repeat costs ~3x more, and deadweight (promo times one-minus-incrementality) shows ~70% of repeat spend is wasted. With more data I'd add iROAS, LTV:CAC, CUPED for variance and uplift modelling to find who's actually persuadable.")

ORDER = [s01, s02, s_map, s03, s_analytics, s05, s_info, s06, s04, s07, s10, s_know, s_persona, s14, s08, s09, s11, s13, s12, s15]

# ENDDEFS
if __name__ == "__main__":
    make_charts()
    prs = build()
    out = os.path.join(os.path.dirname(HERE), "Adrien_Enjalbert_Tripadvisor_AI_Incentives.pptx")
    prs.save(out)
    print("saved", out)
